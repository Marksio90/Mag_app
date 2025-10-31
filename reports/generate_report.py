
from __future__ import annotations
import io, os, json
from pathlib import Path
from typing import Dict, Any, Optional
import pandas as pd

def load_branding(path: str|None) -> dict:
    if not path or not os.path.exists(path): 
        return {}
    try:
        import json
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

# PDF via ReportLab
def build_pdf(report_path: str, logo_path: Optional[str], kpi: dict, figs: Dict[str, bytes], tables: Dict[str, pd.DataFrame], branding_path: Optional[str]=None) -> str:
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet

    theme = load_branding(branding_path).get("palette", {})
    accent = theme.get("accent", "#14b8a6")
    text = theme.get("text", "#222222")
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="TitleBig", fontSize=22, textColor=text, leading=26))
    styles.add(ParagraphStyle(name="H2", fontSize=16, textColor=text, leading=20))
    styles.add(ParagraphStyle(name="H3", fontSize=13, textColor=text, leading=16))

    doc = SimpleDocTemplate(report_path, pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    story = []
    story.append(Paragraph("<b>Zatowarowanie – Executive Report</b>", styles["TitleBig"]))
    if logo_path and os.path.exists(logo_path):
        story.append(Image(logo_path, width=120, height=120))
    story.append(Paragraph(f"<font color='{accent}'>Auto-Story dla Zarządu</font>", styles["H2"]))
    story.append(Spacer(1, 12))

    story.append(Paragraph("<b>KPI danych</b>", styles["H2"]))
    kpi_tbl = pd.DataFrame(list(kpi.items()), columns=["Wskaźnik","Wartość"])
    data = [kpi_tbl.columns.tolist()] + kpi_tbl.values.tolist()
    t = Table(data)
    t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.lightgrey),
                           ("GRID",(0,0),(-1,-1),0.5,colors.grey)]))
    story.append(t)
    story.append(Spacer(1, 12))

    for title, png_bytes in figs.items():
        story.append(Paragraph(f"<b>{title}</b>", styles["H3"]))
        if png_bytes:
            story.append(Image(io.BytesIO(png_bytes), width=520, height=260))
            story.append(Spacer(1, 12))

    for title, df in tables.items():
        story.append(Paragraph(f"<b>{title}</b>", styles["H3"]))
        head = [df.columns.tolist()] + df.head(15).astype(str).values.tolist()
        t = Table(head)
        t.setStyle(TableStyle([("BACKGROUND",(0,0),(-1,0),colors.lightgrey),
                               ("GRID",(0,0),(-1,-1),0.25,colors.grey)]))
        story.append(t)
        story.append(Spacer(1, 12))

    doc.build(story)
    return report_path

# PPTX via python-pptx
def build_pptx(pptx_path: str, logo_path: Optional[str], kpi: dict, figs: Dict[str, bytes], tables: Dict[str, pd.DataFrame], branding_path: Optional[str]=None) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    brand = load_branding(branding_path)
    pal = brand.get("palette", {})
    accent = pal.get("accent", "#14b8a6")
    accent_rgb = RGBColor(int(accent[1:3],16), int(accent[3:5],16), int(accent[5:7],16))

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)).text_frame
    p = tf.paragraphs[0]; p.text = "Zatowarowanie – Executive Report"; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = accent_rgb
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(8), Inches(0.2), height=Inches(1))

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5)).text_frame
    tf.text = "KPI danych"
    y = Inches(1.0)
    for k,v in kpi.items():
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.3)).text_frame
        box.paragraphs[0].text = f"{k}: {v}"
        y += Inches(0.35)

    for title, png_bytes in figs.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4)).text_frame.text = title
        if png_bytes:
            slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(0.5), Inches(1.0), width=Inches(9))

    for title, df in tables.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4)).text_frame.text = title
        y = Inches(1.0)
        head = " | ".join([str(c) for c in df.columns])
        slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.3)).text_frame.text = head
        y += Inches(0.35)
        for _,row in df.head(15).iterrows():
            slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.3)).text_frame.text = " | ".join([str(x) for x in row.values])
            y += Inches(0.3)

    prs.save(pptx_path)
    return pptx_path
