
from __future__ import annotations
import os, io, json
from typing import Dict, Any, List
import plotly.graph_objects as go
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches, Pt

def fig_to_png_bytes(fig: go.Figure) -> bytes:
    # Requires kaleido
    return fig.to_image(format="png", scale=2)

def export_pdf(path: str, title: str, figs: List[go.Figure], tables: List[Dict[str, Any]]):
    c = canvas.Canvas(path, pagesize=A4)
    width, height = A4
    y = height - 50
    c.setFont("Helvetica-Bold", 16)
    c.drawString(40, y, title); y -= 30
    for i, fig in enumerate(figs, start=1):
        try:
            img = fig_to_png_bytes(fig)
            tmp = path + f".tmp{i}.png"
            with open(tmp, "wb") as f: f.write(img)
            c.drawImage(tmp, 40, y-240, width=520, height=240, preserveAspectRatio=True, anchor='n')
            y -= 260
            os.remove(tmp)
        except Exception as e:
            c.setFont("Helvetica", 10); c.drawString(40, y, f"[Błąd renderu wykresu {i}: {e}]"); y -= 14
        if y < 120:
            c.showPage(); y = height - 50
    # tables as JSON blocks
    for t in tables:
        c.setFont("Helvetica-Bold", 12); c.drawString(40, y, t.get("title","Tabela")); y -= 16
        c.setFont("Helvetica", 9)
        txt = json.dumps(t.get("data", {}), ensure_ascii=False, indent=2)
        for line in txt.splitlines():
            c.drawString(40, y, line[:110]); y -= 12
            if y < 80: c.showPage(); y = height - 50
    c.save()

def export_pptx(path: str, title: str, figs: List[go.Figure], tables: List[Dict[str, Any]]):
    prs = Presentation()
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Optimal Inventory Planner — raport"
    for fig in figs:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img = fig_to_png_bytes(fig)
        tmp = path + ".tmp.png"
        with open(tmp, "wb") as f: f.write(img)
        slide.shapes.add_picture(tmp, Inches(0.5), Inches(0.8), width=Inches(9))
        os.remove(tmp)
    for t in tables:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title if slide.shapes.title else None
        if title_shape: title_shape.text = t.get("title","Tabela")
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.text = json.dumps(t.get("data", {}), ensure_ascii=False, indent=2)
    prs.save(path)
